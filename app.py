"""
CIO Digital Flighting Builder
Streamlit app — run with: streamlit run app.py
Deploy to share: push to GitHub, connect to share.streamlit.io
"""

import io
import json
from datetime import date, timedelta

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# ─── Constants ──────────────────────────────────────────────────────────────

COLOR_OPTIONS = {
    "navy":   ("Navy Blue",    "#1C2E5C"),
    "blue":   ("Steel Blue",   "#3B60A8"),
    "teal":   ("Teal",         "#1E6A6E"),
    "olive":  ("Olive/Tan",    "#7A6945"),
    "brown":  ("Brown",        "#96724B"),
    "salmon": ("Coral/Salmon", "#D96B5F"),
    "pink":   ("Hot Pink",     "#D41A82"),
    "tan":    ("Gold/Tan",     "#C0A060"),
}

DEFAULT_CHANNELS = [
    "LinkedIn Static",
    "LinkedIn Video",
    "LinkedIn Carousel",
    "Zoominfo Display",
    "Sponsored Display",
    "Sponsored Newsletter",
    "CIO Q2 Newsletter",
    "Summit Email",
]

# ─── Session state ───────────────────────────────────────────────────────────

def _init():
    if "flights" not in st.session_state:
        st.session_state.flights = []
    if "channels" not in st.session_state:
        st.session_state.channels = list(DEFAULT_CHANNELS)
    if "pptx_bytes" not in st.session_state:
        st.session_state.pptx_bytes = None
    if "edit_idx" not in st.session_state:
        st.session_state.edit_idx = None


def next_flight_id() -> str:
    """Return the smallest unused F-ID (e.g. 'F1', 'F2') based on current flights."""
    used = set()
    for f in st.session_state.flights:
        if isinstance(f.get("id"), str) and f["id"].startswith("F"):
            try:
                used.add(int(f["id"][1:]))
            except ValueError:
                pass
    n = 1
    while n in used:
        n += 1
    return f"F{n}"

_init()

# ─── Date helpers ────────────────────────────────────────────────────────────

def build_week_dates(start: date, n_weeks: int):
    return [start + timedelta(weeks=i) for i in range(n_weeks)]

def week_label(d: date) -> str:
    return f"{d.month}/{d.day}"

def flight_cols(f, week_dates):
    """Return (start_col, end_col_exclusive) for a flight."""
    start_col = end_col = None
    for i, ws in enumerate(week_dates):
        we = ws + timedelta(days=6)
        if f["start_date"] <= we and f["end_date"] >= ws:
            if start_col is None:
                start_col = i
            end_col = i + 1
    return start_col, end_col

def assign_rows(flights, week_dates):
    """Pack flights into rows with no overlap. Returns list of (flight, row_idx)."""
    rows = []
    ordered = sorted(flights, key=lambda f: f["start_date"])
    for fl in ordered:
        sc, ec = flight_cols(fl, week_dates)
        if sc is None:
            continue
        placed = False
        for row in rows:
            ok = True
            for ex in row:
                ex_sc, ex_ec = flight_cols(ex, week_dates)
                if ex_sc < ec and sc < ex_ec:
                    ok = False
                    break
            if ok:
                row.append(fl)
                placed = True
                break
        if not placed:
            rows.append([fl])
    result = []
    for ri, row in enumerate(rows):
        for fl in row:
            result.append((fl, ri))
    return result, len(rows)

def build_channel_cells(flights, week_dates, channels):
    """For each channel, return a list of badge-ID lists (one per week column)."""
    out = {}
    for ch in channels:
        cells = []
        for i, ws in enumerate(week_dates):
            we = ws + timedelta(days=6)
            badges = [
                f["id"] for f in flights
                if ch in f["channels"]
                and f["start_date"] <= we and f["end_date"] >= ws
            ]
            cells.append(badges)
        out[ch] = cells
    return out

# ─── HTML preview ────────────────────────────────────────────────────────────

def build_preview_html(flights, week_dates, current_week_idx,
                       slide_title, subtitle, channels):
    """Render an HTML preview of the slide that mirrors the PPTX layout."""
    n_wk = len(week_dates)
    flight_rows, n_gr = assign_rows(flights, week_dates)
    n_gr = max(n_gr, 1)
    channel_cells = build_channel_cells(flights, week_dates, channels)
    flights_by_id = {f["id"]: f for f in flights}

    css = f"""
    <style>
    .sp {{ font-family: -apple-system, "Segoe UI", Calibri, sans-serif;
           background: white; border: 1px solid #d0cdc6; padding: 14px 16px;
           color: #1a1a2e; max-width: 100%; box-sizing: border-box; }}
    .sp .ttl {{ font-size: 22px; font-weight: 700; line-height: 1.1; }}
    .sp .sub {{ font-size: 12px; color: #6b6560; margin: 2px 0 10px; }}
    .sp .grid {{ display: grid; grid-template-columns: 120px repeat({n_wk}, 1fr);
                 gap: 0; font-size: 9px; }}
    .sp .c {{ border: 0.5px solid #d8d4cd; padding: 2px 4px; box-sizing: border-box;
              min-height: 22px; display: flex; align-items: center;
              overflow: hidden; }}
    .sp .lab {{ background: #eeeae4; font-weight: 600; font-size: 9.5px;
                line-height: 1.15; }}
    .sp .hdr {{ background: #eeeae4; justify-content: center; font-weight: 600;
                font-size: 8.5px; }}
    .sp .hdr.cur {{ background: #1a1a2e; color: white; }}
    .sp .bar {{ color: white; font-weight: 500; font-size: 9px;
                gap: 6px; padding: 2px 6px; white-space: nowrap;
                text-overflow: ellipsis; min-height: 22px; border: none;
                border-right: 1px solid rgba(255,255,255,0.6); }}
    .sp .bar .fid {{ font-weight: 700; padding-right: 6px; flex: 0 0 auto;
                     border-right: 1px solid rgba(255,255,255,0.45); }}
    .sp .bar .lbl {{ overflow: hidden; text-overflow: ellipsis; }}
    .sp .ch {{ padding: 3px; align-content: flex-start; flex-wrap: wrap;
               gap: 2px; min-height: 26px; }}
    .sp .ch.cur {{ background: #ece9f4; }}
    .sp .chb {{ display: inline-block; padding: 1px 3px; font-size: 7.5px;
                font-weight: 700; color: white; border-radius: 2px;
                line-height: 1.3; }}
    .sp .ch-name {{ font-weight: 600; font-size: 9.5px; line-height: 1.15;
                    white-space: pre-line; }}
    .sp .leg {{ display: grid; grid-template-columns: repeat(auto-fill, minmax(190px, 1fr));
                gap: 4px 14px; margin-top: 12px; padding: 8px 10px;
                background: #f7f5f2; border: 1px solid #d8d4cd; font-size: 10px; }}
    .sp .leg-ttl {{ grid-column: 1 / -1; font-size: 9px; font-weight: 700;
                    letter-spacing: 0.5px; color: #6b6560; }}
    .sp .leg-item {{ display: flex; align-items: center; gap: 6px; }}
    .sp .leg-sw {{ width: 14px; height: 12px; flex: 0 0 14px; }}
    .sp .leg-id {{ font-weight: 700; }}
    .sp .empty {{ color: #9a948b; font-style: italic; padding: 12px;
                  text-align: center; }}
    </style>
    """

    h = [css, '<div class="sp">']
    h.append(f'<div class="ttl">{slide_title}</div>')
    h.append(f'<div class="sub">{subtitle}</div>')

    if not flights:
        h.append('<div class="empty">No flights yet. Add one in the sidebar to see the preview.</div>')
        h.append('</div>')
        return "".join(h)

    h.append('<div class="grid">')

    # ── Header row ───────────────────────────────────────────────────────
    h.append('<div class="c lab">Content Flights</div>')
    for i, d in enumerate(week_dates):
        cls = "c hdr cur" if i == current_week_idx else "c hdr"
        h.append(f'<div class="{cls}">{week_label(d)}</div>')

    # ── Gantt rows ───────────────────────────────────────────────────────
    by_row = {}
    for f, gr in flight_rows:
        s, e = flight_cols(f, week_dates)
        if s is None:
            continue
        by_row.setdefault(gr, []).append((f, s, e))
    for gr in range(n_gr):
        h.append('<div class="c lab"></div>')
        bars = sorted(by_row.get(gr, []), key=lambda x: x[1])
        cur_col = 0
        for f, s, e in bars:
            while cur_col < s:
                h.append('<div class="c"></div>')
                cur_col += 1
            span = max(1, e - s)
            color = COLOR_OPTIONS[f["color"]][1]
            esc_label = (f["label"] or "").replace("<", "&lt;").replace(">", "&gt;")
            h.append(
                f'<div class="c bar" style="background:{color};'
                f' grid-column: span {span}">'
                f'<span class="fid">{f["id"]}</span>'
                f'<span class="lbl">{esc_label}</span>'
                f'</div>'
            )
            cur_col = e
        while cur_col < n_wk:
            h.append('<div class="c"></div>')
            cur_col += 1

    # ── Channel rows ─────────────────────────────────────────────────────
    for ch in channels:
        h.append(f'<div class="c lab"><span class="ch-name">{ch}</span></div>')
        cells = channel_cells.get(ch, [])
        for wk in range(n_wk):
            badges_html = ""
            if wk < len(cells):
                for fid in cells[wk]:
                    f = flights_by_id.get(fid)
                    color = COLOR_OPTIONS[f["color"]][1] if f else "#888"
                    badges_html += f'<span class="chb" style="background:{color}">{fid}</span>'
            cls = "c ch cur" if wk == current_week_idx else "c ch"
            h.append(f'<div class="{cls}">{badges_html}</div>')

    h.append('</div>')  # /grid

    # ── Legend ───────────────────────────────────────────────────────────
    h.append('<div class="leg">')
    h.append('<div class="leg-ttl">FLIGHT LEGEND</div>')
    seen = set()
    for f in flights:
        if f["id"] in seen:
            continue
        seen.add(f["id"])
        color = COLOR_OPTIONS[f["color"]][1]
        esc_label = (f["label"] or "").replace("<", "&lt;").replace(">", "&gt;")
        h.append(
            f'<div class="leg-item">'
            f'<span class="leg-sw" style="background:{color}"></span>'
            f'<span><span class="leg-id">{f["id"]}</span> {esc_label}</span>'
            f'</div>'
        )
    h.append('</div>')

    h.append('</div>')  # /sp
    return "".join(h)


# ─── PPTX generation ────────────────────────────────────────────────────────

def _c(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

_C = dict(
    white="FFFFFF", black="000000", border="CCCAC4",
    rowAlt="F4F2EF", titleTx="1A1A2E", subTxt="6B6560",
    adobe="FF0000", colHL="1A1A2E", weekBg="EEEAE4",
    weekTx="2D2926", legBg="F7F5F2", chanHL="E8E4F2",
    navy="1C2E5C", blue="3B60A8", teal="1E6A6E",
    olive="7A6945", brown="96724B", salmon="D96B5F",
    pink="D41A82", tan="C0A060",
)

SL_W, SL_H = 13.3, 7.5
ML, MR, LABEL_W = 0.12, 0.12, 1.08


def _box(slide, x, y, w, h, fill=None, lc=None, lw=0.5):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = _c(fill)
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = _c(lc)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def _lbl(slide, text, x, y, w, h, fs=9, bold=False, col="000000",
         align="left", valign="middle", fill=None, lc=None, lw=0.5,
         wrap=True, ml=0.04):
    s = _box(slide, x, y, w, h, fill=fill, lc=lc, lw=lw)
    tf = s.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(0.01)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.vertical_anchor = (MSO_ANCHOR.MIDDLE if valign == "middle"
                          else MSO_ANCHOR.BOTTOM if valign == "bottom"
                          else MSO_ANCHOR.TOP)
    for i, line in enumerate((text or "").split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = p.space_before = Pt(0)
        p.alignment = (PP_ALIGN.CENTER if align == "center"
                       else PP_ALIGN.RIGHT if align == "right"
                       else PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(fs)
        r.font.bold = bold
        r.font.color.rgb = _c(col)
    return s


def _hline(slide, x, y, length, col="CCCAC4", lw=0.5):
    cn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x), Inches(y), Inches(x + length), Inches(y))
    cn.line.color.rgb = _c(col)
    cn.line.width = Pt(lw)


def _vline(slide, x, y, height, col="CCCAC4", lw=0.3):
    cn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x), Inches(y), Inches(x), Inches(y + height))
    cn.line.color.rgb = _c(col)
    cn.line.width = Pt(lw)


def generate_pptx(flights: list, week_dates: list, current_week_idx: int,
                  slide_title: str, subtitle: str, channels: list) -> bytes:
    n_wk = len(week_dates)
    col_w = (SL_W - ML - MR - LABEL_W) / n_wk
    col_x0 = ML + LABEL_W

    flight_rows, n_gr = assign_rows(flights, week_dates)
    n_gr = max(n_gr, 1)
    channel_cells = build_channel_cells(flights, week_dates, channels)

    # Vertical layout — adaptive so adding many channels still fits
    T_Y, T_H       = 0.10, 0.40
    HDR_Y, HDR_H   = 0.54, 0.26
    G_Y             = HDR_Y + HDR_H
    G_RH            = max(0.15, min(0.22, 2.8 / n_gr))
    G_BOT           = G_Y + n_gr * G_RH
    CH_Y            = G_BOT + 0.07
    FT_Y            = SL_H - 0.26
    LG_H            = 0.70
    # Reserve space for legend + footer; fit channel rows in what's left
    avail_for_chan  = FT_Y - LG_H - 0.06 - CH_Y
    N_CH            = max(len(channels), 1)
    C_RH            = max(0.20, min(0.34, avail_for_chan / N_CH))
    CH_BOT          = CH_Y + N_CH * C_RH
    LG_Y            = CH_BOT + 0.06

    prs = Presentation()
    prs.slide_width = Inches(SL_W)
    prs.slide_height = Inches(SL_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _c(_C["white"])

    # ── Title ──────────────────────────────────────────────────────────────
    _lbl(slide, slide_title, ML, T_Y, 9, T_H * 0.58,
         fs=26, bold=True, col=_C["titleTx"], valign="bottom", ml=0)
    _lbl(slide, subtitle, ML, T_Y + T_H * 0.55, 9, T_H * 0.42,
         fs=13, col=_C["subTxt"], valign="top", ml=0)
    cur_label = f"Week of {week_label(week_dates[current_week_idx])}  \u25cf  Current" if 0 <= current_week_idx < n_wk else ""
    _lbl(slide, cur_label, 8.0, T_Y, 5.1, T_H,
         fs=9, col=_C["subTxt"], align="right", valign="middle", ml=0)

    # ── Date header ────────────────────────────────────────────────────────
    _box(slide, ML, HDR_Y, SL_W - ML - MR, HDR_H, fill=_C["weekBg"], lc=_C["border"])
    _lbl(slide, "Content Flights", ML, HDR_Y, LABEL_W, HDR_H,
         fs=7.5, bold=True, col=_C["titleTx"], valign="middle", ml=0.06)
    for i, wd in enumerate(week_dates):
        x = col_x0 + i * col_w
        hl = (i == current_week_idx)
        _box(slide, x, HDR_Y, col_w, HDR_H,
             fill=_C["colHL"] if hl else _C["weekBg"], lc=_C["border"], lw=0.4)
        _lbl(slide, week_label(wd), x, HDR_Y, col_w, HDR_H,
             fs=7, bold=hl,
             col=_C["white"] if hl else _C["weekTx"],
             align="center", valign="middle", ml=0)

    # ── Gantt rows ─────────────────────────────────────────────────────────
    for r in range(n_gr):
        ry = G_Y + r * G_RH
        _box(slide, ML, ry, SL_W - ML - MR, G_RH,
             fill=_C["white"] if r % 2 == 0 else _C["rowAlt"],
             lc=_C["border"], lw=0.3)
    for i in range(n_wk + 1):
        _vline(slide, col_x0 + i * col_w, G_Y, n_gr * G_RH)
    if 0 <= current_week_idx < n_wk:
        for r in range(n_gr):
            ry = G_Y + r * G_RH
            _box(slide, col_x0 + current_week_idx * col_w, ry, col_w, G_RH,
                 fill="DDD8EE", lc=None)

    G_PAD = 0.02
    BAR_H = G_RH - 2 * G_PAD
    BADGE_W = 0.23

    for fl, row_idx in flight_rows:
        sc, ec = flight_cols(fl, week_dates)
        if sc is None:
            continue
        bx = col_x0 + sc * col_w
        bw = (ec - sc) * col_w
        by = G_Y + row_idx * G_RH + G_PAD
        clr = _C.get(fl["color"], _C["navy"])

        _box(slide, bx, by, bw, BAR_H, fill=clr, lc="FFFFFF", lw=0.4)
        _box(slide, bx, by, BADGE_W, BAR_H, fill=clr, lc="FFFFFF", lw=0.8)
        _lbl(slide, fl["id"], bx, by, BADGE_W, BAR_H,
             fs=7, bold=True, col=_C["white"], align="center", valign="middle", ml=0)
        if bw > BADGE_W + 0.08:
            _lbl(slide, fl["label"],
                 bx + BADGE_W + 0.04, by, bw - BADGE_W - 0.06, BAR_H,
                 fs=6.5, col=_C["white"], align="left", valign="middle", ml=0, wrap=False)

    # ── Channel table ──────────────────────────────────────────────────────
    _hline(slide, ML, CH_Y - 0.02, SL_W - ML - MR, _C["border"], 1.0)

    BADGE_BW  = 0.125
    BADGE_BH  = 0.112
    BADGE_GAP = 0.010
    MAX_ROW   = max(1, int((col_w - 0.015) / (BADGE_BW + BADGE_GAP)))

    for ri, ch in enumerate(channels):
        ry  = CH_Y + ri * C_RH
        alt = ri % 2 == 1
        _box(slide, ML, ry, SL_W - ML - MR, C_RH,
             fill=_C["rowAlt"] if alt else _C["white"], lc=_C["border"], lw=0.3)
        if 0 <= current_week_idx < n_wk:
            _box(slide, col_x0 + current_week_idx * col_w, ry, col_w, C_RH,
                 fill=_C["chanHL"], lc=None)
        for i in range(n_wk + 1):
            _vline(slide, col_x0 + i * col_w, ry, C_RH)
        _lbl(slide, ch, ML, ry, LABEL_W, C_RH,
             fs=7.5, col=_C["titleTx"], align="left", valign="middle", ml=0.06)

        for ci, badges in enumerate(channel_cells[ch]):
            if not badges:
                continue
            cx = col_x0 + ci * col_w
            n_rows = (len(badges) + MAX_ROW - 1) // MAX_ROW
            total_h = n_rows * BADGE_BH + (n_rows - 1) * 0.008
            start_y = ry + (C_RH - total_h) / 2

            for bi, fid in enumerate(badges):
                ri2 = bi // MAX_ROW
                ci2 = bi % MAX_ROW
                bx  = cx + 0.008 + ci2 * (BADGE_BW + BADGE_GAP)
                bby = start_y + ri2 * (BADGE_BH + 0.008)
                fl_match = next((f for f in flights if f["id"] == fid), None)
                bclr = _C.get(fl_match["color"], _C["navy"]) if fl_match else _C["navy"]

                _box(slide, bx, bby, BADGE_BW, BADGE_BH, fill=bclr, lc="FFFFFF", lw=0.3)
                _lbl(slide, fid, bx, bby, BADGE_BW, BADGE_BH,
                     fs=5, bold=True, col=_C["white"],
                     align="center", valign="middle", ml=0)

    # ── Legend ─────────────────────────────────────────────────────────────
    _box(slide, ML, LG_Y, SL_W - ML - MR, LG_H, fill=_C["legBg"], lc=_C["border"], lw=0.5)
    _lbl(slide, "FLIGHT LEGEND", ML + 0.08, LG_Y + 0.05, 1.4, 0.18,
         fs=6.5, bold=True, col=_C["subTxt"], valign="middle", ml=0)

    items_per_row = 8
    entry_w = (SL_W - ML - MR - 0.16) / items_per_row
    LG_BOX_W, LG_BOX_H = 0.15, 0.13

    for idx, fl in enumerate(flights):
        col_i = idx % items_per_row
        row_i = idx // items_per_row
        lx = ML + 0.08 + col_i * entry_w
        ly = LG_Y + 0.27 + row_i * 0.22
        clr = _C.get(fl["color"], _C["navy"])
        _box(slide, lx, ly, LG_BOX_W, LG_BOX_H, fill=clr, lc="FFFFFF", lw=0.3)
        _lbl(slide, fl["id"], lx, ly, LG_BOX_W, LG_BOX_H,
             fs=5, bold=True, col=_C["white"], align="center", valign="middle", ml=0)
        _lbl(slide, fl["label"], lx + LG_BOX_W + 0.04, ly,
             entry_w - LG_BOX_W - 0.06, LG_BOX_H,
             fs=6, col=_C["titleTx"], align="left", valign="middle", ml=0)

    # ── Footer ─────────────────────────────────────────────────────────────
    _box(slide, 0, FT_Y, SL_W, SL_H - FT_Y, fill="F0EDE8", lc=None)
    _hline(slide, 0, FT_Y, SL_W, _C["border"], 0.5)
    _lbl(slide, "Adobe", ML, FT_Y, 1.2, SL_H - FT_Y,
         fs=16, bold=True, col=_C["adobe"], valign="middle", ml=0)
    _lbl(slide, "\u00a9 2024 Adobe. All Rights Reserved. Adobe Confidential.",
         SL_W - 5.5, FT_Y, 5.5 - MR, SL_H - FT_Y,
         fs=7, col=_C["subTxt"], align="right", valign="middle", ml=0)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─── Streamlit UI ────────────────────────────────────────────────────────────

st.set_page_config(
    page_title="Digital Flighting Builder",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
    /* Compact the top padding */
    .block-container { padding-top: 1.5rem; padding-bottom: 1rem; }
    /* Color dot helper */
    .cdot { display:inline-block; width:12px; height:12px;
            border-radius:3px; margin-right:6px; vertical-align:middle; }
    /* Flight card */
    .flight-card { background:#f8f7f5; border:1px solid #e0ddd8;
                   border-radius:8px; padding:10px 14px; margin-bottom:8px; }
    .flight-card h4 { margin:0 0 4px 0; font-size:15px; }
    .flight-card p  { margin:0; font-size:12px; color:#6b6560; }
    .badge { display:inline-block; color:#fff; font-weight:700;
             font-size:10px; border-radius:3px; padding:1px 5px; margin-right:3px; }
    /* Subtle section headers */
    .section-hdr { font-size:11px; font-weight:700; letter-spacing:1px;
                   color:#9b9590; text-transform:uppercase; margin-bottom:4px; }
</style>
""", unsafe_allow_html=True)

# ─── Sidebar ─────────────────────────────────────────────────────────────────

with st.sidebar:
    st.markdown("## ⚙️ Slide Settings")

    slide_title = st.text_input("Slide title", value="CIO Digital Flighting")
    subtitle    = st.text_input("Subtitle", value="Hero Content")

    # Default the slide's week-range start to the Monday of the current week
    _today = date.today()
    _default_range_start = _today - timedelta(days=_today.weekday())

    def _sync_form_dates():
        """When the slide's week range start changes, update the Add-Flight defaults too."""
        new_start = st.session_state.range_start
        st.session_state.f_start = new_start
        st.session_state.f_end   = new_start + timedelta(weeks=4)

    start_date = st.date_input(
        "Week range start",
        value=_default_range_start,
        key="range_start",
        on_change=_sync_form_dates,
    )
    n_weeks    = st.slider("Number of weeks", min_value=6, max_value=26, value=14)

    week_dates = build_week_dates(start_date, n_weeks)
    today = date.today()
    cur_idx = next(
        (i for i, ws in enumerate(week_dates)
         if ws <= today <= ws + timedelta(days=6)),
        0
    )
    cur_idx = st.slider(
        "Current week column",
        min_value=0, max_value=n_weeks - 1, value=cur_idx,
        format=f"Week %d",
        help="Which column is highlighted as 'current week'",
    )

    st.divider()
    with st.expander("📡 Manage Channels", expanded=False):
        st.caption("Channels appear as rows on the slide and as options when adding a flight.")

        # Current channels with delete buttons
        for i, ch in enumerate(st.session_state.channels):
            ccol1, ccol2 = st.columns([5, 1])
            ccol1.write(f"• {ch}")
            if ccol2.button("✕", key=f"del_ch_{i}", help=f"Remove '{ch}'"):
                # Strip from any flights that reference it
                for f in st.session_state.flights:
                    f["channels"] = [c for c in f["channels"] if c != ch]
                st.session_state.channels.pop(i)
                st.session_state.pptx_bytes = None
                st.rerun()

        with st.form("add_channel_form", clear_on_submit=True):
            new_ch = st.text_input(
                "New channel name",
                placeholder="e.g. Reddit Promoted Post",
                label_visibility="collapsed",
            )
            add_ch = st.form_submit_button("＋ Add channel", use_container_width=True)
        if add_ch:
            name = new_ch.strip()
            if not name:
                st.warning("Enter a channel name.")
            elif name in st.session_state.channels:
                st.warning(f"'{name}' already exists.")
            else:
                st.session_state.channels.append(name)
                st.session_state.pptx_bytes = None
                st.success(f"Added '{name}'.")
                st.rerun()

    st.divider()
    st.markdown("## ✈️ Add New Flight")

    auto_id = next_flight_id()
    with st.form("add_flight_form", clear_on_submit=True):
        # Dynamic key — when auto_id changes (e.g. after a delete), the widget
        # resets to the new default instead of clinging to the previous value.
        flight_id  = st.text_input("Flight ID", value=auto_id, key=f"f_id_{auto_id}")
        flight_lbl = st.text_input("Offer / Content label *", placeholder="e.g. 3/5 AIDT Core")

        color_key = st.selectbox(
            "Color",
            options=list(COLOR_OPTIONS.keys()),
            format_func=lambda k: COLOR_OPTIONS[k][0],
        )

        col1, col2 = st.columns(2)
        with col1:
            f_start = st.date_input("Start date", value=start_date, key="f_start")
        with col2:
            f_end = st.date_input("End date", value=start_date + timedelta(weeks=4), key="f_end")

        channels = st.multiselect(
            "Channels",
            options=st.session_state.channels,
            default=[st.session_state.channels[0]] if st.session_state.channels else [],
        )

        submitted = st.form_submit_button("＋ Add Flight", use_container_width=True, type="primary")

    if submitted:
        if not flight_lbl.strip():
            st.error("Offer label is required.")
        elif f_start > f_end:
            st.error("Start date must be before end date.")
        else:
            st.session_state.flights.append({
                "id":         flight_id.strip() or auto_id,
                "label":      flight_lbl.strip(),
                "color":      color_key,
                "start_date": f_start,
                "end_date":   f_end,
                "channels":   channels,
            })
            st.session_state.pptx_bytes = None  # invalidate
            st.success(f"Added {flight_id}!")
            st.rerun()

    st.divider()
    st.markdown("## 💾 Import / Export")

    export_data = json.dumps(
        {
            "channels": st.session_state.channels,
            "flights": [
                {**f, "start_date": str(f["start_date"]), "end_date": str(f["end_date"])}
                for f in st.session_state.flights
            ],
        },
        indent=2,
    )
    st.download_button(
        "Export flights as JSON",
        data=export_data,
        file_name="flights.json",
        mime="application/json",
        use_container_width=True,
    )

    uploaded = st.file_uploader("Import flights from JSON", type="json", label_visibility="collapsed")
    if uploaded:
        try:
            raw = json.loads(uploaded.read())
            # Back-compat: old exports were a bare list of flights
            if isinstance(raw, list):
                imported_flights = raw
                imported_channels = None
            else:
                imported_flights = raw.get("flights", [])
                imported_channels = raw.get("channels")
            for f in imported_flights:
                f["start_date"] = date.fromisoformat(f["start_date"])
                f["end_date"]   = date.fromisoformat(f["end_date"])
            st.session_state.flights = imported_flights
            if imported_channels:
                st.session_state.channels = imported_channels
            st.session_state.pptx_bytes = None
            st.success(f"Imported {len(imported_flights)} flights.")
            st.rerun()
        except Exception as e:
            st.error(f"Import failed: {e}")


# ─── Main area ───────────────────────────────────────────────────────────────

st.markdown("# 📊 Digital Flighting Builder")
st.caption("Build content flight plans and export a polished PowerPoint slide.")

tab_flights, tab_preview, tab_generate = st.tabs(
    ["✈️  Manage Flights", "👁  Preview", "📥  Generate Slide"]
)

# ── Tab 1: Manage Flights ────────────────────────────────────────────────────
with tab_flights:
    if not st.session_state.flights:
        st.info("No flights yet — add one using the sidebar form.")
    else:
        st.markdown(f"**{len(st.session_state.flights)} flight(s)**")

        # Column headers
        col_ratios = [0.6, 2.5, 1.0, 1.0, 1.0, 2.3, 0.4, 0.4]
        hcols = st.columns(col_ratios)
        for col, header in zip(hcols, ["ID", "Label", "Color", "Start", "End", "Channels", "", ""]):
            col.markdown(f"<div class='section-hdr'>{header}</div>", unsafe_allow_html=True)

        st.divider()

        color_keys = list(COLOR_OPTIONS.keys())

        for i, fl in enumerate(st.session_state.flights):
            # ── EDIT MODE ────────────────────────────────────────────────
            if st.session_state.edit_idx == i:
                with st.container(border=True):
                    st.markdown(f"##### ✏️ Editing **{fl['id']}**")

                    e_label = st.text_input(
                        "Label / Offer name",
                        value=fl["label"],
                        key=f"e_label_{i}",
                    )

                    ec1, ec2, ec3 = st.columns([2, 1, 1])
                    e_color = ec1.selectbox(
                        "Color",
                        options=color_keys,
                        index=color_keys.index(fl["color"]) if fl["color"] in color_keys else 0,
                        format_func=lambda k: COLOR_OPTIONS[k][0],
                        key=f"e_color_{i}",
                    )
                    e_start = ec2.date_input(
                        "Start date",
                        value=fl["start_date"],
                        key=f"e_start_{i}",
                    )
                    e_end = ec3.date_input(
                        "End date",
                        value=fl["end_date"],
                        key=f"e_end_{i}",
                    )

                    e_chans = st.multiselect(
                        "Channels",
                        options=st.session_state.channels,
                        default=[c for c in fl["channels"] if c in st.session_state.channels],
                        key=f"e_chan_{i}",
                    )

                    bsave, bcancel, _ = st.columns([1, 1, 4])
                    if bsave.button("💾  Save", key=f"save_{i}", type="primary", use_container_width=True):
                        if e_end < e_start:
                            st.error("End date must be on or after start date.")
                        elif not (e_label or "").strip():
                            st.error("Label can't be empty.")
                        else:
                            st.session_state.flights[i] = {
                                **fl,
                                "label": e_label.strip(),
                                "color": e_color,
                                "start_date": e_start,
                                "end_date": e_end,
                                "channels": e_chans,
                            }
                            st.session_state.edit_idx = None
                            st.session_state.pptx_bytes = None
                            st.rerun()
                    if bcancel.button("Cancel", key=f"cancel_{i}", use_container_width=True):
                        st.session_state.edit_idx = None
                        st.rerun()
                continue

            # ── DISPLAY MODE ─────────────────────────────────────────────
            cols = st.columns(col_ratios)
            hex_color = COLOR_OPTIONS[fl["color"]][1].lstrip("#")
            badge_html = f"<span class='badge' style='background:#{hex_color}'>{fl['id']}</span>"

            cols[0].markdown(badge_html, unsafe_allow_html=True)
            cols[1].markdown(f"**{fl['label']}**")
            cols[2].markdown(
                f"<span class='cdot' style='background:#{hex_color}'></span>"
                f"{COLOR_OPTIONS[fl['color']][0]}",
                unsafe_allow_html=True,
            )
            cols[3].write(fl["start_date"].strftime("%b %-d"))
            cols[4].write(fl["end_date"].strftime("%b %-d"))
            channel_badges = " ".join(
                f"<span style='font-size:11px;background:#f0ede8;padding:1px 6px;"
                f"border-radius:3px;margin-right:2px'>{ch}</span>"
                for ch in fl["channels"]
            )
            cols[5].markdown(channel_badges or "—", unsafe_allow_html=True)
            if cols[6].button("✏", key=f"edit_{i}", help="Edit this flight"):
                st.session_state.edit_idx = i
                st.rerun()
            if cols[7].button("✕", key=f"del_{i}", help="Delete this flight"):
                st.session_state.flights.pop(i)
                if st.session_state.edit_idx == i:
                    st.session_state.edit_idx = None
                st.session_state.pptx_bytes = None
                st.rerun()

        st.divider()
        if st.button("🗑  Clear all flights", type="secondary"):
            st.session_state.flights = []
            st.session_state.pptx_bytes = None
            st.rerun()


# ── Tab 2: Preview ───────────────────────────────────────────────────────────
with tab_preview:
    st.caption("Live preview — updates as you add or edit flights.")
    preview_html = build_preview_html(
        flights=st.session_state.flights,
        week_dates=week_dates,
        current_week_idx=cur_idx,
        slide_title=slide_title,
        subtitle=subtitle,
        channels=st.session_state.channels,
    )
    # Approximate height based on rows
    n_gr_preview = max(assign_rows(st.session_state.flights, week_dates)[1], 1)
    est_height = 110 + 30 + 24 * n_gr_preview + 32 * len(st.session_state.channels) + 140
    st.components.v1.html(preview_html, height=est_height, scrolling=True)


# ── Tab 3: Generate Slide ────────────────────────────────────────────────────
with tab_generate:
    if not st.session_state.flights:
        st.warning("Add at least one flight before generating.")
    else:
        wk_labels = [week_label(d) for d in week_dates]
        st.markdown(f"""
        **Slide:** {slide_title} — {subtitle}
        **Date range:** {wk_labels[0]} → {wk_labels[-1]}  ({n_weeks} weeks)
        **Current week:** {wk_labels[cur_idx]}
        **Flights:** {len(st.session_state.flights)}
        """)

        # Week range preview
        with st.expander("Week columns preview"):
            st.write(" · ".join(
                f"**{w}**" if i == cur_idx else w
                for i, w in enumerate(wk_labels)
            ))

        st.divider()

        if st.button("🎯  Generate PPTX", type="primary", use_container_width=True):
            with st.spinner("Building slide…"):
                try:
                    st.session_state.pptx_bytes = generate_pptx(
                        flights=st.session_state.flights,
                        week_dates=week_dates,
                        current_week_idx=cur_idx,
                        slide_title=slide_title,
                        subtitle=subtitle,
                        channels=st.session_state.channels,
                    )
                    st.success("Slide ready!")
                except Exception as e:
                    st.error(f"Generation failed: {e}")
                    st.exception(e)

        if st.session_state.pptx_bytes:
            st.download_button(
                label="⬇️  Download PPTX",
                data=st.session_state.pptx_bytes,
                file_name="digital_flighting.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )
