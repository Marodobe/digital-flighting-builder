#!/usr/bin/env python3
"""Generate improved CIO Digital Flighting (Hero Content) slide."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

OUTPUT = "/Users/arielmaroniene/Documents/New project/cio_digital_flighting.pptx"

# ─── Color helpers ─────────────────────────────────────────────────────────
def c(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

C = dict(
    olive   = "7A6945",  # F1, F3
    brown   = "96724B",  # F2
    tan     = "C0A060",  # F4
    navy    = "1C2E5C",  # F5, F6
    blue    = "3B60A8",  # F7, F8
    teal    = "1E6A6E",  # F10
    salmon  = "D96B5F",  # F11
    pink    = "D41A82",  # F14–F18
    f9      = "5A7A3A",  # F9 (olive-green)
    f12     = "5B7EC9",  # F12 (steel blue)
    f13     = "7A5C3A",  # F13 (warm brown)
    white   = "FFFFFF",
    black   = "000000",
    border  = "CCCAC4",
    rowAlt  = "F4F2EF",
    titleTx = "1A1A2E",
    subTxt  = "6B6560",
    adobe   = "FF0000",
    colHL   = "1A1A2E",
    weekBg  = "EEEAE4",
    weekTx  = "2D2926",
    legBg   = "F7F5F2",
    chanHL  = "E8E4F2",
)

FC = dict(
    F1=C["olive"], F2=C["brown"], F3=C["olive"], F4=C["tan"],
    F5=C["navy"],  F6=C["navy"],  F7=C["blue"],  F8=C["blue"],
    F9=C["f9"],    F10=C["teal"], F11=C["salmon"],
    F12=C["f12"],  F13=C["f13"],
    F14=C["pink"], F15=C["pink"], F16=C["pink"],
    F17=C["pink"], F18=C["pink"],
)

# ─── Layout constants ───────────────────────────────────────────────────────
SL_W, SL_H = 13.3, 7.5
ML = 0.12                          # left margin
MR = 0.12                          # right margin
LABEL_W = 1.08
COL_X0  = ML + LABEL_W             # 1.20"
N_WK    = 14
COL_W   = (SL_W - ML - MR - LABEL_W) / N_WK   # ≈ 0.855"

T_Y   = 0.10;  T_H   = 0.40
HDR_Y = 0.54;  HDR_H = 0.26
G_Y   = HDR_Y + HDR_H              # 0.80  gantt top
G_RH  = 0.195                      # gantt row height
N_GR  = 12                         # gantt rows
G_BOT = G_Y + N_GR * G_RH          # 0.80 + 2.34 = 3.14
CH_Y  = G_BOT + 0.07               # 3.21  channel top
C_RH  = 0.295                      # channel row height
N_CH  = 8
CH_BOT= CH_Y + N_CH * C_RH         # 3.21 + 2.36 = 5.57
LG_Y  = CH_BOT + 0.06              # 5.63  legend top
LG_H  = 0.72
FT_Y  = SL_H - 0.26                # 7.24

CUR_WK = 6  # 4/13 column

WEEKS = [
    "3/5","3/9","3/16","3/23","3/30","4/6",
    "4/13","4/20","4/27","5/4","5/11","5/18","5/25","6/1 Q3\u2192",
]

# ─── Flight data ────────────────────────────────────────────────────────────
# (id, color_key, gantt_row, start_col, end_col_excl, bar_label)
FLIGHTS = [
    ("F1",  "olive",  0,  0,  6, '3/5: \u201cLIVE from MWC\u201d video'),
    ("F3",  "olive",  0,  6, 14, "4/11: Trust By Design ODW"),
    ("F4",  "tan",    1,  6, 14, "4/15: Webinar Soundbites"),
    ("F2",  "brown",  2,  2, 14, "3/9: TVW in-feed video ad"),
    ("F5",  "navy",   3,  0,  6, "3/5 AIDT Core"),
    ("F6",  "navy",   3,  6, 14, "4/9 AIDT Core  CIO Fact Sheet"),
    ("F7",  "blue",   4,  0,  6, "3/5 AI Inflection Point"),
    ("F8",  "blue",   4,  6, 14, "4/20 AI Inflection Point Chapter 2: Crossing the AI Adoption Gap"),
    ("F10", "teal",   5,  7, 14, "IDC CMO/CIO  (CIO derivative)"),
    ("F11", "salmon", 6,  0, 14, "CIO Persona Page"),
    ("F14", "pink",   7,  6, 14, "(TBD START DATE) CXO Maturity Index (Summit)"),
    ("F15", "pink",   8,  6, 14, "(TBD START DATE) Top IT Post Summit sessions (Summit)"),
    ("F16", "pink",   9,  6, 14, "(TBD START DATE) Top IT summit recap cutdowns (Summit)"),
    ("F17", "pink",  10,  6, 14, "Partner announcement blogs (Summit)"),
    ("F18", "pink",  11,  6, 14, "Workday / intuit / ulta blogs (Summit)"),
]

# ─── Channel table data ─────────────────────────────────────────────────────
# Each row: channel name + 14-cell list of badge IDs per week column
CHAN_DATA = [
    ("LinkedIn\nStatic", [
        ["F5","F7","F11"],           # 3/5
        ["F5","F7","F11"],           # 3/9
        ["F5","F7","F11"],           # 3/16
        ["F5","F7","F11","F12","F13"],# 3/23
        ["F5","F7","F11","F12","F13"],# 3/30
        ["F5","F7","F11","F12","F13"],# 4/6
        ["F6","F8","F9","F11","F12","F13"],# 4/13 ← current
        ["F6","F8","F9","F11","F4"], # 4/20
        ["F6","F8","F9","F11","F4"], # 4/27
        ["F6","F8","F10","F11","F4"],# 5/4
        ["F6","F8","F10","F11","F4"],# 5/11
        ["F6","F8","F10","F11","F4"],# 5/18
        ["F6","F8","F10","F11","F4"],# 5/25
        ["F6","F8","F10","F11","F4"],# 6/1
    ]),
    ("LinkedIn\nVideo", [
        ["F1","F2"], ["F1","F2"], ["F1","F2"], ["F1","F2"],
        ["F1","F2","F12"],
        ["F2","F3","F12"],
        ["F2","F3","F12"],
        ["F2","F4","F15"],["F2","F4","F15"],["F2","F4","F15"],
        ["F2","F4","F15"],["F2","F4","F15"],["F2","F4","F15"],
        ["F4","F15"],
    ]),
    ("LinkedIn\nCarousel", [
        ["F5"],["F5"],["F5"],["F5"],["F5"],
        ["F5","F9"],
        ["F6","F9"],
        ["F6","F9","F17","F18"],
        ["F6","F9","F17","F18"],
        ["F6","F10","F15","F17","F18"],
        ["F6","F10","F15","F17","F18"],
        ["F6","F10","F15","F17","F18"],
        ["F6","F10","F15","F17","F18"],
        ["F6","F10","F15","F17","F18"],
    ]),
    ("Zoominfo\nDisplay", [
        ["F11"],["F11"],["F11"],["F11"],
        ["F5","F7","F11"],
        ["F5","F7","F9","F11"],
        ["F6","F7","F8","F9","F11"],
        ["F6","F8","F9","F11","F17","F18"],
        ["F6","F8","F10","F11","F17","F18"],
        ["F6","F8","F10","F11","F17","F18"],
        ["F6","F8","F10","F11","F17","F18"],
        ["F6","F8","F10","F11","F17","F18"],
        ["F6","F8","F10","F11","F17","F18"],
        ["F6","F8","F10","F11","F17","F18"],
    ]),
    ("Sponsored\nDisplay", [
        ["F11"],["F11"],["F11"],["F11"],
        ["F5","F7","F9","F11"],
        ["F5","F7","F9","F11"],
        ["F6","F8","F9","F11"],
        ["F6","F8","F9","F11"],
        ["F6","F8","F10","F11"],
        ["F6","F8","F10","F11"],
        ["F6","F8","F10","F11"],
        ["F6","F8","F10","F11"],
        ["F6","F8","F10","F11"],
        ["F6","F8","F10","F11"],
    ]),
    ("Sponsored\nNewsletter", [
        [],[],[],[],[],[],
        ["F6","F8"],
        [],[],[],[],[],[],[],
    ]),
    ("CIO Q2\nNewsletter", [
        [],[],[],[],[],[],[],[],
        [],["F6","F10","F11","F14"],
        [],[],[],[],
    ]),
    ("Summit Email #1\n(4/23)", [
        [],[],[],[],[],[],
        ["F17"],
        [],[],[],[],[],[],[],
    ]),
]

# ─── Drawing helpers ────────────────────────────────────────────────────────
def box(slide, x, y, w, h, fill=None, lc=None, lw=0.5):
    """Add a filled/bordered rectangle via textbox."""
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = c(fill)
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = c(lc)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def label(slide, text, x, y, w, h, fs=9, bold=False, color="000000",
          align="left", valign="middle", fill=None, lc=None, lw=0.5,
          wrap=True, italic=False, ml=0.04, mt=0):
    """Add text in a rectangle."""
    s = box(slide, x, y, w, h, fill=fill, lc=lc, lw=lw)
    tf = s.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left  = Inches(ml)
    tf.margin_right = Inches(0.01)
    tf.margin_top   = Inches(mt)
    tf.margin_bottom= Inches(0)
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = text.split("\n") if text else []
    for i, line_txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.alignment = (PP_ALIGN.CENTER if align == "center"
                       else PP_ALIGN.RIGHT if align == "right"
                       else PP_ALIGN.LEFT)
        r = p.add_run()
        r.text = line_txt
        r.font.size  = Pt(fs)
        r.font.bold  = bold
        r.font.italic= italic
        r.font.color.rgb = c(color)
    return s

def hline(slide, x, y, length, color="CCCAC4", lw=0.5):
    """Horizontal connector line."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x), Inches(y), Inches(x + length), Inches(y))
    conn.line.color.rgb = c(color)
    conn.line.width = Pt(lw)
    return conn

def vline(slide, x, y, height, color="CCCAC4", lw=0.5):
    """Vertical connector line."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x), Inches(y), Inches(x), Inches(y + height))
    conn.line.color.rgb = c(color)
    conn.line.width = Pt(lw)
    return conn

# ─── Build slide ────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(SL_W)
prs.slide_height = Inches(SL_H)

blank_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(blank_layout)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = c(C["white"])

# ── Title ──────────────────────────────────────────────────────────────────
label(slide, "CIO Digital Flighting", ML, T_Y, 8, T_H * 0.58,
      fs=26, bold=True, color=C["titleTx"], valign="bottom", ml=0)
label(slide, "Hero Content", ML, T_Y + T_H * 0.55, 8, T_H * 0.42,
      fs=13, color=C["subTxt"], valign="top", ml=0)
label(slide, "Week of 4/13  \u25cf  Current", 8.0, T_Y, 5.1, T_H,
      fs=9, color=C["subTxt"], align="right", valign="middle", ml=0)

# ── Date header row ────────────────────────────────────────────────────────
box(slide, ML, HDR_Y, SL_W - ML - MR, HDR_H, fill=C["weekBg"], lc=C["border"])
label(slide, "Content Flights", ML, HDR_Y, LABEL_W, HDR_H,
      fs=7.5, bold=True, color=C["titleTx"], valign="middle", ml=0.06)

for i, wk in enumerate(WEEKS):
    x = COL_X0 + i * COL_W
    hl = (i == CUR_WK)
    box(slide, x, HDR_Y, COL_W, HDR_H,
        fill=C["colHL"] if hl else C["weekBg"], lc=C["border"], lw=0.4)
    label(slide, wk, x, HDR_Y, COL_W, HDR_H,
          fs=7, bold=hl, color=C["white"] if hl else C["weekTx"],
          align="center", valign="middle", ml=0)

# ── Gantt section ──────────────────────────────────────────────────────────
# Row backgrounds
for r in range(N_GR):
    ry = G_Y + r * G_RH
    fill = C["white"] if r % 2 == 0 else C["rowAlt"]
    box(slide, ML, ry, SL_W - ML - MR, G_RH, fill=fill, lc=C["border"], lw=0.3)

# Vertical column separators
for i in range(N_WK + 1):
    x = COL_X0 + i * COL_W
    vline(slide, x, G_Y, N_GR * G_RH, C["border"], 0.3)

# Current week subtle column tint in gantt
for r in range(N_GR):
    ry = G_Y + r * G_RH
    box(slide, COL_X0 + CUR_WK * COL_W, ry, COL_W, G_RH,
        fill="DDD8EE", lc=None)

# Flight bars
G_PAD = 0.022
BAR_H = G_RH - 2 * G_PAD
BADGE_W = 0.23

for fid, ck, row, start, end, lbl in FLIGHTS:
    bx  = COL_X0 + start * COL_W
    bw  = (end - start) * COL_W
    by  = G_Y + row * G_RH + G_PAD
    clr = C[ck]

    # Main bar
    box(slide, bx, by, bw, BAR_H, fill=clr, lc="FFFFFF", lw=0.4)

    # Flight ID badge — same bar color, white border separates it visually
    box(slide, bx, by, BADGE_W, BAR_H, fill=clr, lc="FFFFFF", lw=0.8)

    label(slide, fid, bx, by, BADGE_W, BAR_H,
          fs=7, bold=True, color=C["white"], align="center", valign="middle", ml=0)

    # Content label
    if bw > BADGE_W + 0.08:
        label(slide, lbl,
              bx + BADGE_W + 0.04, by, bw - BADGE_W - 0.06, BAR_H,
              fs=6.5, color=C["white"], align="left", valign="middle", ml=0, wrap=False)

# ── Channel section ────────────────────────────────────────────────────────
hline(slide, ML, CH_Y - 0.02, SL_W - ML - MR, C["border"], 1.0)

BADGE_BW  = 0.125
BADGE_BH  = 0.112
BADGE_GAP = 0.010
MAX_ROW   = int((COL_W - 0.015) / (BADGE_BW + BADGE_GAP))  # ≈ 6

for ri, (chan_name, cells) in enumerate(CHAN_DATA):
    ry    = CH_Y + ri * C_RH
    alt   = (ri % 2 == 1)
    bg    = C["rowAlt"] if alt else C["white"]

    # Row background
    box(slide, ML, ry, SL_W - ML - MR, C_RH, fill=bg, lc=C["border"], lw=0.3)

    # Current week column tint
    box(slide, COL_X0 + CUR_WK * COL_W, ry, COL_W, C_RH, fill=C["chanHL"], lc=None)

    # Column separators
    for i in range(N_WK + 1):
        x = COL_X0 + i * COL_W
        vline(slide, x, ry, C_RH, C["border"], 0.3)

    # Channel name
    label(slide, chan_name, ML, ry, LABEL_W, C_RH,
          fs=7.5, color=C["titleTx"], align="left", valign="middle", ml=0.06)

    # Badges
    for ci, badges in enumerate(cells):
        if not badges:
            continue
        cx = COL_X0 + ci * COL_W
        n_rows = (len(badges) + MAX_ROW - 1) // MAX_ROW
        total_h = n_rows * BADGE_BH + (n_rows - 1) * 0.008
        start_y = ry + (C_RH - total_h) / 2

        for bi, fid in enumerate(badges):
            row_i = bi // MAX_ROW
            col_i = bi %  MAX_ROW
            bx = cx + 0.008 + col_i * (BADGE_BW + BADGE_GAP)
            bby = start_y + row_i * (BADGE_BH + 0.008)
            bclr = FC.get(fid, C["black"])

            box(slide, bx, bby, BADGE_BW, BADGE_BH, fill=bclr, lc="FFFFFF", lw=0.3)
            label(slide, fid, bx, bby, BADGE_BW, BADGE_BH,
                  fs=5, bold=True, color=C["white"],
                  align="center", valign="middle", ml=0)

# ── Legend ─────────────────────────────────────────────────────────────────
box(slide, ML, LG_Y, SL_W - ML - MR, LG_H, fill=C["legBg"], lc=C["border"], lw=0.5)

label(slide, "FLIGHT LEGEND", ML + 0.08, LG_Y + 0.05, 1.4, 0.18,
      fs=6.5, bold=True, color=C["subTxt"], valign="middle", ml=0)

LEGEND_ITEMS = [
    ("F1",  "olive",  "LIVE from MWC Video"),
    ("F2",  "brown",  "TVW In-Feed Video Ad"),
    ("F3",  "olive",  "Trust By Design ODW"),
    ("F4",  "tan",    "Webinar Soundbites"),
    ("F5",  "navy",   "AIDT Core"),
    ("F6",  "navy",   "AIDT Core CIO Fact Sheet"),
    ("F7",  "blue",   "AI Inflection Point"),
    ("F8",  "blue",   "AI Inflection Point Ch. 2"),
    ("F10", "teal",   "IDC CMO/CIO"),
    ("F11", "salmon", "CIO Persona Page"),
    ("F14", "pink",   "CXO Maturity Index (Summit)"),
    ("F15", "pink",   "Top IT Post Summit Sessions"),
    ("F16", "pink",   "IT Summit Recap Cutdowns"),
    ("F17", "pink",   "Partner Announcement Blogs"),
    ("F18", "pink",   "Workday/Intuit/Ulta Blogs"),
]

ITEMS_PER_ROW = 8
ENTRY_W = (SL_W - ML - MR - 0.16) / ITEMS_PER_ROW   # ≈ 1.62"
LG_BOX_W = 0.15
LG_BOX_H = 0.13
LG_TXT_W = ENTRY_W - LG_BOX_W - 0.06

for idx, (fid, ck, name) in enumerate(LEGEND_ITEMS):
    col = idx % ITEMS_PER_ROW
    row = idx // ITEMS_PER_ROW
    lx  = ML + 0.08 + col * ENTRY_W
    ly  = LG_Y + 0.27 + row * 0.22

    box(slide, lx, ly, LG_BOX_W, LG_BOX_H, fill=C[ck], lc="FFFFFF", lw=0.3)
    label(slide, fid, lx, ly, LG_BOX_W, LG_BOX_H,
          fs=5, bold=True, color=C["white"], align="center", valign="middle", ml=0)
    label(slide, name, lx + LG_BOX_W + 0.04, ly, LG_TXT_W, LG_BOX_H,
          fs=6, color=C["titleTx"], align="left", valign="middle", ml=0)

# ── Footer ─────────────────────────────────────────────────────────────────
box(slide, 0, FT_Y, SL_W, SL_H - FT_Y, fill="F0EDE8", lc=None)
hline(slide, 0, FT_Y, SL_W, C["border"], 0.5)

label(slide, "Adobe", ML, FT_Y, 1.2, SL_H - FT_Y,
      fs=16, bold=True, color=C["adobe"], valign="middle", ml=0)
label(slide, "\u00a9 2024 Adobe. All Rights Reserved. Adobe Confidential.",
      SL_W - 5.5, FT_Y, 5.5 - MR, SL_H - FT_Y,
      fs=7, color=C["subTxt"], align="right", valign="middle", ml=0)

# ── Save ───────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
